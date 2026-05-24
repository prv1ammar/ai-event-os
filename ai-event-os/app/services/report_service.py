"""
app/services/report_service.py
────────────────────────────────
Post-event report generation:
  • 8-page PDF     (ReportLab)
  • 7-slide PPTX   (python-pptx)
  • Multi-sheet XLSX (openpyxl)

PDF pages
─────────
  1. Cover
  2. Executive Summary (KPI cards)
  3. Visitor Analytics
  4. Exhibitor Performance
  5. Leads & ROI
  6. Marketing Performance
  7. AI Recommendations
  8. Appendix (raw data tables)

PPTX slides (dark navy #1a1a2e + purple #7c3aed theme)
────────────────────────────────────────────────────────
  1. Title — event name + headline KPIs
  2. Attendance Overview
  3. Exhibitor Performance
  4. Leads & Business
  5. Marketing
  6. Financial Results & ROI
  7. Recommendations & Next Edition
"""
from __future__ import annotations

import io
import uuid
from datetime import datetime, timezone
from typing import Any

from sqlalchemy import func, select
from sqlalchemy.ext.asyncio import AsyncSession

# ─────────────────────────────────────────────────────────────────────────────
# Colour palette (shared by PDF and PPTX)
# ─────────────────────────────────────────────────────────────────────────────
NAVY   = (26, 26, 46)       # #1a1a2e
PURPLE = (124, 58, 237)     # #7c3aed
GOLD   = (245, 158, 11)     # amber accent
WHITE  = (255, 255, 255)
LIGHT  = (241, 245, 249)    # light grey
DARK   = (15, 23, 42)       # almost black


# ═════════════════════════════════════════════════════════════════════════════
# PDF — ReportLab
# ═════════════════════════════════════════════════════════════════════════════

def generate_post_event_pdf(event: Any, stats: dict) -> bytes:
    """
    Generate a full 8-page post-event PDF report.

    Parameters
    ──────────
    event : SQLAlchemy Event ORM object (or dict with name, venue, start_date…)
    stats : dict returned by analytics_service.get_dashboard_kpis()

    Returns
    ───────
    bytes  — PDF binary ready to stream as a FileResponse.
    """
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm, mm
    from reportlab.platypus import (
        BaseDocTemplate, Frame, PageBreak, PageTemplate,
        Paragraph, Spacer, Table, TableStyle,
    )
    from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT

    # ── Colour helpers ─────────────────────────────────────────────────────────
    def rgb(r, g, b):
        return colors.Color(r / 255, g / 255, b / 255)

    C_NAVY   = rgb(*NAVY)
    C_PURPLE = rgb(*PURPLE)
    C_GOLD   = rgb(*GOLD)
    C_LIGHT  = rgb(*LIGHT)
    C_WHITE  = rgb(*WHITE)

    # ── Styles ─────────────────────────────────────────────────────────────────
    styles = getSampleStyleSheet()

    title_style = ParagraphStyle(
        "Title",
        parent=styles["Title"],
        fontSize=28, textColor=C_WHITE, spaceAfter=8, alignment=TA_CENTER,
    )
    h1_style = ParagraphStyle(
        "H1",
        parent=styles["Heading1"],
        fontSize=18, textColor=C_NAVY, spaceAfter=6,
    )
    h2_style = ParagraphStyle(
        "H2",
        parent=styles["Heading2"],
        fontSize=13, textColor=C_PURPLE, spaceAfter=4,
    )
    body_style = ParagraphStyle(
        "Body",
        parent=styles["Normal"],
        fontSize=10, leading=14, spaceAfter=4,
    )
    kpi_style = ParagraphStyle(
        "KPI",
        parent=styles["Normal"],
        fontSize=24, textColor=C_PURPLE, alignment=TA_CENTER, spaceAfter=2,
    )
    caption_style = ParagraphStyle(
        "Caption",
        parent=styles["Normal"],
        fontSize=9, textColor=colors.grey, alignment=TA_CENTER,
    )
    white_body = ParagraphStyle(
        "WhiteBody",
        parent=styles["Normal"],
        fontSize=11, textColor=C_WHITE, alignment=TA_CENTER,
    )

    # ── Event metadata ─────────────────────────────────────────────────────────
    event_name  = getattr(event, "name", None)  or stats.get("event_name", "Event")
    event_venue = getattr(event, "venue", None) or "N/A"
    event_city  = getattr(event, "city", None)  or ""
    start_date  = getattr(event, "start_date", None)
    end_date    = getattr(event, "end_date", None)

    date_str = ""
    if start_date:
        date_str = str(start_date)
        if end_date:
            date_str += f" – {end_date}"

    # ── KPI values ─────────────────────────────────────────────────────────────
    total_visitors  = stats.get("total_visitors", 0)
    total_exhibitors= stats.get("total_exhibitors", 0)
    total_leads     = stats.get("total_leads", 0)
    qualified_leads = stats.get("qualified_leads", 0)
    meetings        = stats.get("meetings_scheduled", 0)
    revenue         = stats.get("total_revenue_mad", 0)
    expenses        = stats.get("total_expenses_mad", 0)
    roi             = stats.get("roi_percent", 0)
    occupancy       = stats.get("occupancy_rate", 0)
    avg_score       = stats.get("avg_lead_score", 0)
    conv_rate       = round(qualified_leads / total_leads * 100, 1) if total_leads > 0 else 0

    buffer = io.BytesIO()
    W, H = A4

    # ── Document & frames ──────────────────────────────────────────────────────
    margin = 2 * cm
    frame = Frame(margin, margin, W - 2 * margin, H - 2 * margin, id="main")

    def cover_bg(canvas, doc):
        canvas.saveState()
        canvas.setFillColor(C_NAVY)
        canvas.rect(0, 0, W, H, fill=1, stroke=0)
        # Purple accent bar
        canvas.setFillColor(C_PURPLE)
        canvas.rect(0, H * 0.35, W, 4, fill=1, stroke=0)
        canvas.restoreState()

    def normal_page(canvas, doc):
        canvas.saveState()
        # Header line
        canvas.setFillColor(C_PURPLE)
        canvas.rect(0, H - 1.2 * cm, W, 0.4 * cm, fill=1, stroke=0)
        # Footer
        canvas.setFillColor(C_NAVY)
        canvas.setFont("Helvetica", 8)
        canvas.drawCentredString(W / 2, 0.7 * cm,
                                 f"AI EVENT OS — {event_name} — Confidential")
        canvas.drawRightString(W - margin, 0.7 * cm, f"Page {doc.page}")
        canvas.restoreState()

    cover_template  = PageTemplate(id="cover",  frames=[frame], onPage=cover_bg)
    normal_template = PageTemplate(id="normal", frames=[frame], onPage=normal_page)

    doc = BaseDocTemplate(
        buffer, pagesize=A4,
        pageTemplates=[cover_template, normal_template],
        leftMargin=margin, rightMargin=margin,
        topMargin=margin + 0.5 * cm, bottomMargin=margin,
    )

    story = []

    # ── PAGE 1: Cover ──────────────────────────────────────────────────────────
    story.append(Spacer(1, 5 * cm))
    story.append(Paragraph("AI EVENT OS", white_body))
    story.append(Spacer(1, 0.5 * cm))
    story.append(Paragraph(event_name, title_style))
    story.append(Spacer(1, 0.3 * cm))
    story.append(Paragraph("Post-Event Report", white_body))
    story.append(Spacer(1, 0.5 * cm))
    story.append(Paragraph(
        f"{event_venue}{', ' + event_city if event_city else ''}",
        ParagraphStyle("sub", parent=white_body, fontSize=10, textColor=colors.lightgrey)
    ))
    story.append(Paragraph(
        date_str,
        ParagraphStyle("sub", parent=white_body, fontSize=10, textColor=colors.lightgrey)
    ))
    story.append(Spacer(1, 2 * cm))
    story.append(Paragraph(
        f"Generated: {datetime.now(timezone.utc).strftime('%Y-%m-%d')}",
        ParagraphStyle("gen", parent=white_body, fontSize=9, textColor=colors.grey)
    ))
    story.append(PageBreak())

    # ── PAGE 2: Executive Summary ──────────────────────────────────────────────
    story.append(Paragraph("Executive Summary", h1_style))
    story.append(Spacer(1, 0.4 * cm))

    def kpi_table(pairs):
        data = []
        row_label = []
        row_value = []
        for label, value in pairs:
            row_label.append(Paragraph(str(label), caption_style))
            row_value.append(Paragraph(str(value), kpi_style))
        data = [row_value, row_label]
        t = Table(data, colWidths=[(W - 4 * cm) / len(pairs)] * len(pairs))
        t.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), C_LIGHT),
            ("BOX", (0, 0), (-1, -1), 0.5, colors.lightgrey),
            ("INNERGRID", (0, 0), (-1, -1), 0.25, colors.lightgrey),
            ("TOPPADDING", (0, 0), (-1, -1), 10),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 10),
        ]))
        return t

    story.append(kpi_table([
        ("Total Visitors",  f"{total_visitors:,}"),
        ("Exhibitors",      f"{total_exhibitors:,}"),
        ("Total Leads",     f"{total_leads:,}"),
        ("B2B Meetings",    f"{meetings:,}"),
    ]))
    story.append(Spacer(1, 0.6 * cm))
    story.append(kpi_table([
        ("Revenue (MAD)",   f"{revenue:,.0f}"),
        ("Expenses (MAD)",  f"{expenses:,.0f}"),
        ("ROI",             f"{roi:.1f}%"),
        ("Occupancy",       f"{occupancy:.1f}%"),
    ]))
    story.append(Spacer(1, 0.8 * cm))

    # Performance vs objectives table
    story.append(Paragraph("Performance vs Objectives", h2_style))
    perf_data = [
        ["KPI", "Target", "Actual", "Status"],
        ["Visitors",   "1 000",  f"{total_visitors:,}",  "✓" if total_visitors >= 1000 else "✗"],
        ["Leads",      "500",    f"{total_leads:,}",      "✓" if total_leads >= 500 else "✗"],
        ["Conversion", "15%",    f"{conv_rate:.1f}%",     "✓" if conv_rate >= 15 else "✗"],
        ["ROI",        "100%",   f"{roi:.1f}%",           "✓" if roi >= 100 else "✗"],
        ["Occupancy",  "80%",    f"{occupancy:.1f}%",     "✓" if occupancy >= 80 else "✗"],
    ]
    perf_t = Table(perf_data, colWidths=[5 * cm, 3 * cm, 3 * cm, 2 * cm])
    perf_t.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), C_NAVY),
        ("TEXTCOLOR",  (0, 0), (-1, 0), C_WHITE),
        ("FONTNAME",   (0, 0), (-1, 0), "Helvetica-Bold"),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [C_WHITE, C_LIGHT]),
        ("GRID", (0, 0), (-1, -1), 0.5, colors.lightgrey),
        ("FONTSIZE", (0, 0), (-1, -1), 9),
        ("TOPPADDING", (0, 0), (-1, -1), 6),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
    ]))
    story.append(perf_t)
    story.append(PageBreak())

    # ── PAGE 3: Visitor Analytics ──────────────────────────────────────────────
    story.append(Paragraph("Visitor Analytics", h1_style))
    story.append(Spacer(1, 0.4 * cm))

    v_breakdown = stats.get("visitor_type_breakdown", {})
    if v_breakdown:
        story.append(Paragraph("Visitor Type Breakdown", h2_style))
        vt_data = [["Type", "Count", "% Total"]] + [
            [k.capitalize(), str(v), f"{v / max(sum(v_breakdown.values()), 1) * 100:.1f}%"]
            for k, v in v_breakdown.items()
        ]
        vt_t = Table(vt_data, colWidths=[5 * cm, 4 * cm, 4 * cm])
        vt_t.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), C_PURPLE),
            ("TEXTCOLOR",  (0, 0), (-1, 0), C_WHITE),
            ("FONTNAME",   (0, 0), (-1, 0), "Helvetica-Bold"),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [C_WHITE, C_LIGHT]),
            ("GRID", (0, 0), (-1, -1), 0.5, colors.lightgrey),
            ("FONTSIZE", (0, 0), (-1, -1), 9),
            ("TOPPADDING", (0, 0), (-1, -1), 6),
        ]))
        story.append(vt_t)
        story.append(Spacer(1, 0.6 * cm))

    story.append(Paragraph("Key Figures", h2_style))
    story.append(Paragraph(
        f"<b>{total_visitors:,}</b> visitors registered · "
        f"<b>{stats.get('entries_today', 0):,}</b> entries recorded · "
        f"Capacity: <b>{stats.get('capacity', 0):,}</b>",
        body_style,
    ))
    story.append(PageBreak())

    # ── PAGE 4: Exhibitor Performance ─────────────────────────────────────────
    story.append(Paragraph("Exhibitor Performance", h1_style))
    story.append(Spacer(1, 0.4 * cm))
    story.append(kpi_table([
        ("Total Exhibitors",  f"{total_exhibitors:,}"),
        ("Booths Reserved",   f"{stats.get('reserved_booths', 0):,}"),
        ("Occupancy Rate",    f"{occupancy:.1f}%"),
        ("Avg Lead Score",    f"{avg_score:.1f} / 100"),
    ]))
    story.append(Spacer(1, 0.8 * cm))
    story.append(Paragraph(
        "Booth occupancy measures how many booths were reserved or occupied out "
        "of the total available floor space. The average AI lead score reflects "
        "the quality of business interest generated across all exhibitor stands.",
        body_style,
    ))
    story.append(PageBreak())

    # ── PAGE 5: Leads & ROI ────────────────────────────────────────────────────
    story.append(Paragraph("Leads & Business ROI", h1_style))
    story.append(Spacer(1, 0.4 * cm))

    # Funnel table
    story.append(Paragraph("Lead Funnel", h2_style))
    funnel_data = [
        ["Stage",         "Count",                  "Conversion"],
        ["New Leads",     f"{total_leads:,}",        "100%"],
        ["Qualified",     f"{qualified_leads:,}",    f"{conv_rate:.1f}%"],
        ["Meetings",      f"{meetings:,}",           f"{round(meetings / max(total_leads, 1) * 100, 1):.1f}%"],
    ]
    f_t = Table(funnel_data, colWidths=[6 * cm, 4 * cm, 4 * cm])
    f_t.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), C_NAVY),
        ("TEXTCOLOR",  (0, 0), (-1, 0), C_WHITE),
        ("FONTNAME",   (0, 0), (-1, 0), "Helvetica-Bold"),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [C_WHITE, C_LIGHT]),
        ("GRID", (0, 0), (-1, -1), 0.5, colors.lightgrey),
        ("FONTSIZE", (0, 0), (-1, -1), 10),
        ("TOPPADDING", (0, 0), (-1, -1), 8),
    ]))
    story.append(f_t)
    story.append(Spacer(1, 0.8 * cm))
    story.append(kpi_table([
        ("Revenue MAD",  f"{revenue:,.0f}"),
        ("Expenses MAD", f"{expenses:,.0f}"),
        ("Net Profit",   f"{revenue - expenses:,.0f} MAD"),
        ("ROI",          f"{roi:.1f}%"),
    ]))
    story.append(PageBreak())

    # ── PAGE 6: Marketing Performance ─────────────────────────────────────────
    story.append(Paragraph("Marketing Performance", h1_style))
    story.append(Spacer(1, 0.4 * cm))

    traffic = stats.get("traffic_sources", [])
    if traffic:
        story.append(Paragraph("Traffic Sources", h2_style))
        ts_data = [["Channel", "Leads Generated", "Sent Count"]] + [
            [r["source"].capitalize(), str(r["leads_generated"]), str(r["sent_count"])]
            for r in traffic
        ]
        ts_t = Table(ts_data, colWidths=[5 * cm, 5 * cm, 4 * cm])
        ts_t.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), C_PURPLE),
            ("TEXTCOLOR",  (0, 0), (-1, 0), C_WHITE),
            ("FONTNAME",   (0, 0), (-1, 0), "Helvetica-Bold"),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [C_WHITE, C_LIGHT]),
            ("GRID", (0, 0), (-1, -1), 0.5, colors.lightgrey),
            ("FONTSIZE", (0, 0), (-1, -1), 9),
            ("TOPPADDING", (0, 0), (-1, -1), 6),
        ]))
        story.append(ts_t)
    else:
        story.append(Paragraph(
            "No sent campaign data available for this event.", body_style
        ))
    story.append(PageBreak())

    # ── PAGE 7: AI Recommendations ────────────────────────────────────────────
    story.append(Paragraph("AI-Generated Recommendations", h1_style))
    story.append(Spacer(1, 0.4 * cm))

    recommendations = stats.get("insights", [
        f"Booth occupancy was {occupancy:.1f}% — target ≥ 80% for next edition.",
        f"Lead conversion rate was {conv_rate:.1f}% — industry benchmark is 15–20%.",
        f"ROI was {roi:.1f}% — continue optimising the exhibitor package mix.",
        "Introduce a 48-hour pre-event B2B matching reminder campaign.",
        "Add a VIP lounge session to increase high-value visitor engagement.",
    ])

    for i, rec in enumerate(recommendations[:6], 1):
        story.append(Paragraph(f"<b>{i}.</b> {rec}", body_style))
        story.append(Spacer(1, 0.2 * cm))

    story.append(PageBreak())

    # ── PAGE 8: Appendix ──────────────────────────────────────────────────────
    story.append(Paragraph("Appendix — Key Data", h1_style))
    story.append(Spacer(1, 0.4 * cm))

    appendix_data = [
        ["Metric", "Value"],
        ["Event Name",      event_name],
        ["Venue",           f"{event_venue}, {event_city}"],
        ["Dates",           date_str],
        ["Total Visitors",  f"{total_visitors:,}"],
        ["Total Exhibitors",f"{total_exhibitors:,}"],
        ["Total Leads",     f"{total_leads:,}"],
        ["Qualified Leads", f"{qualified_leads:,}"],
        ["Conversion Rate", f"{conv_rate:.1f}%"],
        ["B2B Meetings",    f"{meetings:,}"],
        ["Revenue (MAD)",   f"{revenue:,.0f}"],
        ["Expenses (MAD)",  f"{expenses:,.0f}"],
        ["Net Profit (MAD)",f"{revenue - expenses:,.0f}"],
        ["ROI",             f"{roi:.1f}%"],
        ["Booth Occupancy", f"{occupancy:.1f}%"],
        ["Avg Lead Score",  f"{avg_score:.1f} / 100"],
        ["Report Date",     datetime.now(timezone.utc).strftime("%Y-%m-%d")],
    ]

    app_t = Table(appendix_data, colWidths=[7 * cm, 7 * cm])
    app_t.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), C_NAVY),
        ("TEXTCOLOR",  (0, 0), (-1, 0), C_WHITE),
        ("FONTNAME",   (0, 0), (-1, 0), "Helvetica-Bold"),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [C_WHITE, C_LIGHT]),
        ("GRID", (0, 0), (-1, -1), 0.5, colors.lightgrey),
        ("FONTSIZE", (0, 0), (-1, -1), 9),
        ("TOPPADDING", (0, 0), (-1, -1), 5),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
    ]))
    story.append(app_t)

    # ── Build PDF ──────────────────────────────────────────────────────────────
    doc.build(story)
    return buffer.getvalue()


# ═════════════════════════════════════════════════════════════════════════════
# PPTX — python-pptx
# ═════════════════════════════════════════════════════════════════════════════

def generate_executive_pptx(event: Any, stats: dict) -> bytes:
    """
    Generate a 7-slide executive PowerPoint presentation.

    Theme: dark navy (#1a1a2e) background + purple (#7c3aed) accent.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    def rgb_color(r, g, b):
        return RGBColor(r, g, b)

    C_NAVY_RGB   = rgb_color(*NAVY)
    C_PURPLE_RGB = rgb_color(*PURPLE)
    C_GOLD_RGB   = rgb_color(*GOLD)
    C_WHITE_RGB  = rgb_color(*WHITE)
    C_LIGHT_RGB  = rgb_color(200, 210, 230)

    # ── Presentation setup ─────────────────────────────────────────────────────
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # completely blank

    # ── Event metadata ─────────────────────────────────────────────────────────
    event_name  = getattr(event, "name", None)  or stats.get("event_name", "Event")
    event_venue = getattr(event, "venue", None) or "N/A"
    event_city  = getattr(event, "city", None)  or ""
    start_date  = getattr(event, "start_date", None)
    end_date    = getattr(event, "end_date", None)
    date_str    = f"{start_date} – {end_date}" if start_date and end_date else str(start_date or "")

    # ── KPI aliases ────────────────────────────────────────────────────────────
    total_visitors  = stats.get("total_visitors", 0)
    total_exhibitors= stats.get("total_exhibitors", 0)
    total_leads     = stats.get("total_leads", 0)
    qualified_leads = stats.get("qualified_leads", 0)
    meetings        = stats.get("meetings_scheduled", 0)
    revenue         = stats.get("total_revenue_mad", 0)
    expenses        = stats.get("total_expenses_mad", 0)
    roi             = stats.get("roi_percent", 0)
    occupancy       = stats.get("occupancy_rate", 0)
    avg_score       = stats.get("avg_lead_score", 0)
    conv_rate       = round(qualified_leads / max(total_leads, 1) * 100, 1)

    # ── Helper functions ───────────────────────────────────────────────────────
    W = prs.slide_width
    H = prs.slide_height

    def fill_slide_bg(slide, color_rgb):
        """Fill the slide background with a solid colour."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color_rgb

    def add_rect(slide, left, top, width, height, fill_rgb, line_rgb=None):
        from pptx.util import Emu
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            left, top, width, height,
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
        if line_rgb:
            shape.line.color.rgb = line_rgb
        else:
            shape.line.fill.background()
        return shape

    def add_text_box(slide, text, left, top, width, height,
                     font_size=18, bold=False, color=None, align=PP_ALIGN.LEFT,
                     italic=False):
        from pptx.util import Pt
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text)
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        if color:
            run.font.color.rgb = color
        return txBox

    def kpi_card(slide, left, top, width, height, value, label,
                 bg_rgb=None, val_color=None, lbl_color=None):
        """Add a KPI card with big number + label."""
        bg = bg_rgb or C_WHITE_RGB
        add_rect(slide, left, top, width, height, bg)
        val_c = val_color or C_PURPLE_RGB
        lbl_c = lbl_color or C_NAVY_RGB
        add_text_box(slide, value,
                     left, top + Inches(0.15), width, Inches(0.65),
                     font_size=32, bold=True, color=val_c, align=PP_ALIGN.CENTER)
        add_text_box(slide, label,
                     left, top + Inches(0.8), width, Inches(0.35),
                     font_size=11, bold=False, color=lbl_c, align=PP_ALIGN.CENTER)

    # ═══════════════════════════════════════════════════════════════════════════
    # SLIDE 1: Title
    # ═══════════════════════════════════════════════════════════════════════════
    slide1 = prs.slides.add_slide(blank_layout)
    fill_slide_bg(slide1, C_NAVY_RGB)

    # Purple accent strip
    add_rect(slide1, 0, Inches(2.8), W, Inches(0.06), C_PURPLE_RGB)

    add_text_box(slide1, "AI EVENT OS",
                 Inches(0.5), Inches(0.4), W - Inches(1), Inches(0.5),
                 font_size=14, bold=True, color=C_PURPLE_RGB, align=PP_ALIGN.CENTER)

    add_text_box(slide1, event_name,
                 Inches(0.5), Inches(1.0), W - Inches(1), Inches(1.5),
                 font_size=40, bold=True, color=C_WHITE_RGB, align=PP_ALIGN.CENTER)

    add_text_box(slide1, "Post-Event Executive Report",
                 Inches(0.5), Inches(2.5), W - Inches(1), Inches(0.5),
                 font_size=16, color=C_LIGHT_RGB, align=PP_ALIGN.CENTER)

    add_text_box(slide1, f"{event_venue}  ·  {date_str}",
                 Inches(0.5), Inches(3.1), W - Inches(1), Inches(0.4),
                 font_size=12, color=C_LIGHT_RGB, align=PP_ALIGN.CENTER)

    # Big KPI row
    kw = Inches(2.5)
    gap = Inches(0.18)
    kpi_top = Inches(4.0)
    kh = Inches(1.3)

    cards = [
        (f"{total_visitors:,}", "Visitors"),
        (f"{total_exhibitors:,}", "Exhibitors"),
        (f"{total_leads:,}", "Total Leads"),
        (f"{roi:.0f}%", "ROI"),
        (f"{occupancy:.0f}%", "Occupancy"),
    ]
    start_left = (W - (kw + gap) * len(cards) + gap) / 2
    for i, (val, lbl) in enumerate(cards):
        kpi_card(slide1,
                 start_left + i * (kw + gap), kpi_top,
                 kw, kh, val, lbl,
                 bg_rgb=RGBColor(35, 35, 60),
                 val_color=C_GOLD_RGB,
                 lbl_color=C_LIGHT_RGB)

    add_text_box(slide1,
                 f"Generated: {datetime.now(timezone.utc).strftime('%Y-%m-%d')}",
                 Inches(0.5), Inches(6.9), W - Inches(1), Inches(0.3),
                 font_size=9, color=RGBColor(100, 100, 130), align=PP_ALIGN.CENTER)

    # ═══════════════════════════════════════════════════════════════════════════
    # SLIDE 2: Attendance Overview
    # ═══════════════════════════════════════════════════════════════════════════
    slide2 = prs.slides.add_slide(blank_layout)
    fill_slide_bg(slide2, C_NAVY_RGB)
    add_rect(slide2, 0, 0, W, Inches(1.0), C_PURPLE_RGB)

    add_text_box(slide2, "Attendance Overview",
                 Inches(0.3), Inches(0.15), W - Inches(0.6), Inches(0.7),
                 font_size=28, bold=True, color=C_WHITE_RGB)

    attendance_kpis = [
        (f"{total_visitors:,}", "Total Registered"),
        (f"{stats.get('entries_today', 0):,}", "Entries Today"),
        (f"{stats.get('entries_total', 0):,}", "Total Entries"),
        (f"{stats.get('capacity', 0):,}", "Capacity"),
    ]
    kw2 = Inches(2.8)
    start2 = Inches(0.5)
    top2 = Inches(1.3)
    for i, (val, lbl) in enumerate(attendance_kpis):
        kpi_card(slide2,
                 start2 + i * (kw2 + Inches(0.25)), top2,
                 kw2, Inches(1.5), val, lbl,
                 bg_rgb=RGBColor(35, 35, 60),
                 val_color=C_GOLD_RGB, lbl_color=C_LIGHT_RGB)

    # Visitor type breakdown
    v_breakdown = stats.get("visitor_type_breakdown", {})
    if v_breakdown:
        add_text_box(slide2, "Visitor Type Breakdown",
                     Inches(0.5), Inches(3.2), Inches(6), Inches(0.4),
                     font_size=14, bold=True, color=C_PURPLE_RGB)
        y_off = Inches(3.8)
        for vtype, count in list(v_breakdown.items())[:6]:
            total_v = max(sum(v_breakdown.values()), 1)
            pct = count / total_v * 100
            add_text_box(slide2,
                         f"{vtype.capitalize()}: {count:,}  ({pct:.1f}%)",
                         Inches(0.5), y_off, Inches(5), Inches(0.35),
                         font_size=11, color=C_LIGHT_RGB)
            y_off += Inches(0.4)

    # ═══════════════════════════════════════════════════════════════════════════
    # SLIDE 3: Exhibitor Performance
    # ═══════════════════════════════════════════════════════════════════════════
    slide3 = prs.slides.add_slide(blank_layout)
    fill_slide_bg(slide3, C_NAVY_RGB)
    add_rect(slide3, 0, 0, W, Inches(1.0), C_PURPLE_RGB)
    add_text_box(slide3, "Exhibitor Performance",
                 Inches(0.3), Inches(0.15), W - Inches(0.6), Inches(0.7),
                 font_size=28, bold=True, color=C_WHITE_RGB)

    ex_kpis = [
        (f"{total_exhibitors:,}", "Total Exhibitors"),
        (f"{stats.get('reserved_booths', 0):,}", "Booths Reserved"),
        (f"{occupancy:.1f}%", "Occupancy Rate"),
        (f"{avg_score:.1f}", "Avg Lead Score"),
    ]
    kw3 = Inches(2.8)
    for i, (val, lbl) in enumerate(ex_kpis):
        kpi_card(slide3,
                 Inches(0.5) + i * (kw3 + Inches(0.25)), Inches(1.3),
                 kw3, Inches(1.5), val, lbl,
                 bg_rgb=RGBColor(35, 35, 60),
                 val_color=C_GOLD_RGB, lbl_color=C_LIGHT_RGB)

    add_text_box(slide3,
                 "Booth occupancy measures stand utilisation against total available capacity. "
                 "Average lead score (0–100) reflects the AI-assessed quality of business "
                 "interactions generated across all exhibitor stands.",
                 Inches(0.5), Inches(3.2), W - Inches(1), Inches(1.2),
                 font_size=12, color=C_LIGHT_RGB)

    # ═══════════════════════════════════════════════════════════════════════════
    # SLIDE 4: Leads & Business
    # ═══════════════════════════════════════════════════════════════════════════
    slide4 = prs.slides.add_slide(blank_layout)
    fill_slide_bg(slide4, C_NAVY_RGB)
    add_rect(slide4, 0, 0, W, Inches(1.0), C_PURPLE_RGB)
    add_text_box(slide4, "Leads & Business Development",
                 Inches(0.3), Inches(0.15), W - Inches(0.6), Inches(0.7),
                 font_size=28, bold=True, color=C_WHITE_RGB)

    lead_kpis = [
        (f"{total_leads:,}", "Total Leads"),
        (f"{qualified_leads:,}", "Qualified"),
        (f"{conv_rate:.1f}%", "Conversion"),
        (f"{meetings:,}", "B2B Meetings"),
    ]
    kw4 = Inches(2.8)
    for i, (val, lbl) in enumerate(lead_kpis):
        kpi_card(slide4,
                 Inches(0.5) + i * (kw4 + Inches(0.25)), Inches(1.3),
                 kw4, Inches(1.5), val, lbl,
                 bg_rgb=RGBColor(35, 35, 60),
                 val_color=C_GOLD_RGB, lbl_color=C_LIGHT_RGB)

    # Funnel text
    add_text_box(slide4, "Lead Funnel",
                 Inches(0.5), Inches(3.2), Inches(5), Inches(0.4),
                 font_size=14, bold=True, color=C_PURPLE_RGB)

    funnel_items = [
        f"New Leads:       {total_leads:,}",
        f"Qualified:       {qualified_leads:,}  ({conv_rate:.1f}%)",
        f"B2B Meetings:    {meetings:,}  ({round(meetings / max(total_leads, 1) * 100, 1):.1f}%)",
    ]
    y_f = Inches(3.8)
    for item in funnel_items:
        add_text_box(slide4, item, Inches(0.5), y_f, Inches(7), Inches(0.35),
                     font_size=12, color=C_LIGHT_RGB)
        y_f += Inches(0.4)

    # ═══════════════════════════════════════════════════════════════════════════
    # SLIDE 5: Marketing
    # ═══════════════════════════════════════════════════════════════════════════
    slide5 = prs.slides.add_slide(blank_layout)
    fill_slide_bg(slide5, C_NAVY_RGB)
    add_rect(slide5, 0, 0, W, Inches(1.0), C_PURPLE_RGB)
    add_text_box(slide5, "Marketing Performance",
                 Inches(0.3), Inches(0.15), W - Inches(0.6), Inches(0.7),
                 font_size=28, bold=True, color=C_WHITE_RGB)

    traffic = stats.get("traffic_sources", [])
    if traffic:
        add_text_box(slide5, "Campaign Channel Performance",
                     Inches(0.5), Inches(1.2), W - Inches(1), Inches(0.4),
                     font_size=14, bold=True, color=C_PURPLE_RGB)
        y_tr = Inches(1.8)
        for row in traffic:
            add_text_box(
                slide5,
                f"{row['source'].capitalize():12s}  ·  "
                f"Leads: {row['leads_generated']:,}  ·  Sent: {row['sent_count']:,}",
                Inches(0.5), y_tr, W - Inches(1), Inches(0.35),
                font_size=11, color=C_LIGHT_RGB,
            )
            y_tr += Inches(0.4)
    else:
        add_text_box(slide5, "No campaign data available for this event.",
                     Inches(0.5), Inches(2.0), W - Inches(1), Inches(0.5),
                     font_size=13, color=C_LIGHT_RGB)

    # ═══════════════════════════════════════════════════════════════════════════
    # SLIDE 6: Financial Results & ROI
    # ═══════════════════════════════════════════════════════════════════════════
    slide6 = prs.slides.add_slide(blank_layout)
    fill_slide_bg(slide6, C_NAVY_RGB)
    add_rect(slide6, 0, 0, W, Inches(1.0), C_PURPLE_RGB)
    add_text_box(slide6, "Financial Results & ROI",
                 Inches(0.3), Inches(0.15), W - Inches(0.6), Inches(0.7),
                 font_size=28, bold=True, color=C_WHITE_RGB)

    fin_kpis = [
        (f"{revenue:,.0f}", "Revenue (MAD)"),
        (f"{expenses:,.0f}", "Expenses (MAD)"),
        (f"{revenue - expenses:,.0f}", "Net Profit (MAD)"),
        (f"{roi:.1f}%", "ROI"),
    ]
    kw6 = Inches(2.8)
    for i, (val, lbl) in enumerate(fin_kpis):
        kpi_card(slide6,
                 Inches(0.5) + i * (kw6 + Inches(0.25)), Inches(1.3),
                 kw6, Inches(1.5), val, lbl,
                 bg_rgb=RGBColor(35, 35, 60),
                 val_color=C_GOLD_RGB, lbl_color=C_LIGHT_RGB)

    budget_mad = stats.get("total_budget_mad", 0)
    if budget_mad:
        budget_var = round((expenses - budget_mad) / budget_mad * 100, 1) if budget_mad else 0
        add_text_box(slide6,
                     f"Budget: {budget_mad:,.0f} MAD  ·  "
                     f"Variance: {budget_var:+.1f}%",
                     Inches(0.5), Inches(3.2), W - Inches(1), Inches(0.4),
                     font_size=13, color=C_LIGHT_RGB)

    # ═══════════════════════════════════════════════════════════════════════════
    # SLIDE 7: Recommendations & Next Edition
    # ═══════════════════════════════════════════════════════════════════════════
    slide7 = prs.slides.add_slide(blank_layout)
    fill_slide_bg(slide7, C_NAVY_RGB)
    add_rect(slide7, 0, 0, W, Inches(1.0), C_PURPLE_RGB)
    add_text_box(slide7, "Recommendations & Next Edition",
                 Inches(0.3), Inches(0.15), W - Inches(0.6), Inches(0.7),
                 font_size=28, bold=True, color=C_WHITE_RGB)

    recs = stats.get("insights", [
        f"Target booth occupancy ≥ 85% (current: {occupancy:.1f}%).",
        f"Improve lead conversion from {conv_rate:.1f}% to 20%+ with better matching.",
        "Launch B2B matching app 2 weeks before next edition.",
        "Add VIP networking dinners to increase high-value visitor retention.",
        "Optimise marketing mix — focus on top-performing channels.",
    ])

    y_rec = Inches(1.3)
    for i, rec in enumerate(recs[:6], 1):
        # Bullet circle
        add_rect(slide7,
                 Inches(0.4), y_rec + Inches(0.05),
                 Inches(0.25), Inches(0.25), C_PURPLE_RGB)
        add_text_box(slide7, f"{i}.", Inches(0.4), y_rec + Inches(0.02),
                     Inches(0.25), Inches(0.3),
                     font_size=11, bold=True, color=C_WHITE_RGB, align=PP_ALIGN.CENTER)
        add_text_box(slide7, rec,
                     Inches(0.8), y_rec, W - Inches(1.1), Inches(0.4),
                     font_size=12, color=C_LIGHT_RGB)
        y_rec += Inches(0.55)

    add_text_box(slide7,
                 f"AI EVENT OS  ·  {event_name}  ·  Post-Event Report",
                 Inches(0.5), Inches(7.0), W - Inches(1), Inches(0.3),
                 font_size=9, color=RGBColor(100, 100, 130), align=PP_ALIGN.CENTER)

    # ── Serialise ──────────────────────────────────────────────────────────────
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ═════════════════════════════════════════════════════════════════════════════
# XLSX — openpyxl
# ═════════════════════════════════════════════════════════════════════════════

def generate_event_excel(event: Any, stats: dict) -> bytes:
    """
    Generate an Excel workbook with 4 sheets:
      1. KPI Summary
      2. Visitor Breakdown
      3. Financial Summary
      4. Lead Funnel
    """
    from openpyxl import Workbook
    from openpyxl.styles import (
        Alignment, Border, Font, PatternFill, Side,
    )

    wb = Workbook()

    def hex_fill(hex_str: str) -> PatternFill:
        return PatternFill("solid", fgColor=hex_str)

    navy_fill   = hex_fill("1A1A2E")
    purple_fill = hex_fill("7C3AED")
    light_fill  = hex_fill("F1F5F9")
    gold_fill   = hex_fill("F59E0B")

    def style_header(cell, bg_fill=navy_fill):
        cell.fill = bg_fill
        cell.font = Font(bold=True, color="FFFFFF", size=11)
        cell.alignment = Alignment(horizontal="center", vertical="center")

    def style_kpi_value(cell):
        cell.font = Font(bold=True, size=14, color="7C3AED")
        cell.alignment = Alignment(horizontal="center", vertical="center")

    # ── Sheet 1: KPI Summary ───────────────────────────────────────────────────
    ws1 = wb.active
    ws1.title = "KPI Summary"

    kpi_rows = [
        ("Event", stats.get("event_name", "")),
        ("Status", stats.get("event_status", "")),
        ("Total Visitors", stats.get("total_visitors", 0)),
        ("Total Entries", stats.get("total_entries", 0)),
        ("Entries Today", stats.get("entries_today", 0)),
        ("Total Exhibitors", stats.get("total_exhibitors", 0)),
        ("Total Leads", stats.get("total_leads", 0)),
        ("Qualified Leads", stats.get("qualified_leads", 0)),
        ("B2B Meetings", stats.get("meetings_scheduled", 0)),
        ("Revenue (MAD)", stats.get("total_revenue_mad", 0)),
        ("Expenses (MAD)", stats.get("total_expenses_mad", 0)),
        ("Net Profit (MAD)", stats.get("total_revenue_mad", 0) - stats.get("total_expenses_mad", 0)),
        ("ROI %", stats.get("roi_percent", 0)),
        ("Booth Occupancy %", stats.get("occupancy_rate", 0)),
        ("Avg Lead Score", stats.get("avg_lead_score", 0)),
        ("Capacity", stats.get("capacity", 0)),
    ]

    ws1.append(["Metric", "Value"])
    style_header(ws1.cell(1, 1))
    style_header(ws1.cell(1, 2))
    ws1.column_dimensions["A"].width = 25
    ws1.column_dimensions["B"].width = 20

    for i, (label, value) in enumerate(kpi_rows, 2):
        ws1.cell(i, 1).value = label
        ws1.cell(i, 2).value = value
        if i % 2 == 0:
            ws1.cell(i, 1).fill = light_fill
            ws1.cell(i, 2).fill = light_fill

    # ── Sheet 2: Visitor Breakdown ────────────────────────────────────────────
    ws2 = wb.create_sheet("Visitor Breakdown")
    ws2.append(["Type", "Count", "% Total"])
    style_header(ws2.cell(1, 1), purple_fill)
    style_header(ws2.cell(1, 2), purple_fill)
    style_header(ws2.cell(1, 3), purple_fill)
    ws2.column_dimensions["A"].width = 20
    ws2.column_dimensions["B"].width = 15
    ws2.column_dimensions["C"].width = 15

    v_breakdown = stats.get("visitor_type_breakdown", {})
    total_v = max(sum(v_breakdown.values()), 1)
    for i, (vtype, count) in enumerate(v_breakdown.items(), 2):
        ws2.cell(i, 1).value = vtype.capitalize()
        ws2.cell(i, 2).value = count
        ws2.cell(i, 3).value = round(count / total_v * 100, 1)
        if i % 2 == 0:
            for col in range(1, 4):
                ws2.cell(i, col).fill = light_fill

    # ── Sheet 3: Financial Summary ────────────────────────────────────────────
    ws3 = wb.create_sheet("Financial Summary")
    ws3.append(["Category", "Amount (MAD)"])
    style_header(ws3.cell(1, 1), navy_fill)
    style_header(ws3.cell(1, 2), navy_fill)
    ws3.column_dimensions["A"].width = 30
    ws3.column_dimensions["B"].width = 20

    fin_rows: list[tuple] = [
        ("Total Revenue", stats.get("total_revenue_mad", 0)),
        ("Total Expenses", stats.get("total_expenses_mad", 0)),
        ("Net Profit", stats.get("total_revenue_mad", 0) - stats.get("total_expenses_mad", 0)),
        ("Budget", stats.get("total_budget_mad", 0)),
        ("ROI %", stats.get("roi_percent", 0)),
    ]

    rev_by_src = stats.get("revenue_by_source", {})
    for src, amt in rev_by_src.items():
        fin_rows.append((f"  Revenue — {src}", amt))

    exp_by_cat = stats.get("expenses_by_category", {})
    for cat, amt in exp_by_cat.items():
        fin_rows.append((f"  Expense — {cat}", amt))

    for i, (label, value) in enumerate(fin_rows, 2):
        ws3.cell(i, 1).value = label
        ws3.cell(i, 2).value = round(float(value), 2)
        if i % 2 == 0:
            ws3.cell(i, 1).fill = light_fill
            ws3.cell(i, 2).fill = light_fill

    # ── Sheet 4: Lead Funnel ──────────────────────────────────────────────────
    ws4 = wb.create_sheet("Lead Funnel")
    ws4.append(["Stage", "Count", "Conversion %"])
    style_header(ws4.cell(1, 1), purple_fill)
    style_header(ws4.cell(1, 2), purple_fill)
    style_header(ws4.cell(1, 3), purple_fill)
    ws4.column_dimensions["A"].width = 25
    ws4.column_dimensions["B"].width = 15
    ws4.column_dimensions["C"].width = 18

    total_leads = max(stats.get("total_leads", 0), 1)
    qualified   = stats.get("qualified_leads", 0)
    meetings    = stats.get("meetings_scheduled", 0)

    funnel = [
        ("New Leads",     stats.get("total_leads", 0), 100.0),
        ("Qualified",     qualified, round(qualified / total_leads * 100, 1)),
        ("B2B Meetings",  meetings,  round(meetings  / total_leads * 100, 1)),
    ]
    for i, (stage, count, pct) in enumerate(funnel, 2):
        ws4.cell(i, 1).value = stage
        ws4.cell(i, 2).value = count
        ws4.cell(i, 3).value = pct
        if i % 2 == 0:
            for col in range(1, 4):
                ws4.cell(i, col).fill = light_fill

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ═════════════════════════════════════════════════════════════════════════════
# DB-backed report builder (called by router)
# ═════════════════════════════════════════════════════════════════════════════

async def build_event_report_data(
    db: AsyncSession, event_id: uuid.UUID
) -> tuple[Any, dict]:
    """
    Fetch the Event object + full stats dict for report generation.
    Returns (event, stats_dict).
    """
    from app.models.event import Event as EventModel
    from app.services.analytics_service import (
        get_dashboard_kpis,
        get_traffic_sources,
        get_visitor_type_breakdown,
    )

    result = await db.execute(
        select(EventModel).where(EventModel.id == event_id)
    )
    event = result.scalar_one_or_none()
    if event is None:
        from fastapi import HTTPException, status
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail=f"Event {event_id} not found",
        )

    stats = await get_dashboard_kpis(db, event_id)
    traffic = await get_traffic_sources(db, event_id)
    v_breakdown = await get_visitor_type_breakdown(db, event_id)

    stats["traffic_sources"] = traffic
    stats["visitor_type_breakdown"] = v_breakdown

    return event, stats
